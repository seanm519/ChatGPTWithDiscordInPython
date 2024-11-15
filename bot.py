import os
import discord
from discord.ext import commands
from dotenv import load_dotenv
import asyncio
import aiohttp
from io import BytesIO
import pdfplumber  
#import docx  
from docx import Document
from pptx import Presentation  
import difflib
from discord.ui import View, Button  # Importing View and Button
from collections import deque
from PIL import Image
import pytesseract
import io
from cachetools import TTLCache
import pytesseract

load_dotenv()  # Load environment variables from .env file

DISCORD_TOKEN = os.getenv('DISCORD_TOKEN')
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')

intents = discord.Intents.default()
intents.message_content = True  # Allow bot to read message content
intents.messages = True  # Allow bot to listen for messages


bot = commands.Bot(command_prefix='/', intents=intents)

# Name of the admin role
ADMIN_ROLE_NAME = "Professor"

# In-memory caches
lectures_cache = {}
interactions_cache = []

# Queue to handle incoming messages
question_queue = deque()

# Semaphore to limit the number of concurrent API requests
semaphore = asyncio.Semaphore(3)  # Limit to 3 concurrent requests

# Function to check if the user has the administrator role
async def is_admin(interaction):
    admin_role = discord.utils.get(interaction.guild.roles, name=ADMIN_ROLE_NAME)
    if admin_role in interaction.user.roles:
        return True
    else:
        await interaction.response.send_message("You do not have the required role to use this command.", ephemeral=True)
        return False


# Function to create the admin role if it doesn't exist
async def create_admin_role(guild):
    existing_role = discord.utils.get(guild.roles, name=ADMIN_ROLE_NAME)
    if not existing_role:
        try:
            await guild.create_role(
                name=ADMIN_ROLE_NAME,
                permissions=discord.Permissions(administrator=True),
                mentionable=True
            )
            print(f"Created '{ADMIN_ROLE_NAME}' role in '{guild.name}'.")
        except discord.Forbidden:
            print(f"Bot does not have permissions to create roles in '{guild.name}'.")

# Function to summarize text in manageable chunks
async def summarize_text(text: str) -> str:
    print("Starting text summarization...")
    chunk_size = 1999  # Example chunk size
    text_chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
    print(f"Text split into {len(text_chunks)} chunks")

    summaries = []
    async with aiohttp.ClientSession() as session:
        for i, chunk in enumerate(text_chunks):
            try:
                print(f"Processing chunk {i+1}/{len(text_chunks)}")

                async with session.post(
                    'https://api.openai.com/v1/chat/completions',
                    headers={
                        'Authorization': f'Bearer {OPENAI_API_KEY}',
                        'Content-Type': 'application/json'
                    },
                    json={
                        'model': 'gpt-4',
                        'messages': [{'role': 'user', 'content': f"Summarize the following text:\n{chunk}"}]
                    }
                ) as response:
                    response_json = await response.json()
                    
                    if response.status == 200:
                        summary_text = response_json['choices'][0]['message']['content']
                        summaries.append(summary_text)
                        print(f"Chunk {i+1} summarized successfully")
                    else:
                        print(f"Error summarizing chunk {i+1}: {response_json}")
                        summaries.append("Error in summarizing this chunk.")

            except Exception as e:
                print(f"Exception while processing chunk {i+1}: {e}")
                summaries.append("Error in summarizing this chunk due to an exception.")
                
    final_summary = " ".join(summaries)
    print("Summarization complete.")
    return final_summary

async def read_file_content(attachment: discord.Attachment) -> str:
    try:
        file_content = await attachment.read()
        file_like_object = BytesIO(file_content)

        # PDF Handling with incremental reading
        if attachment.filename.endswith(".pdf"):
            try:
                text = ""
                with pdfplumber.open(file_like_object) as pdf:
                    for page in pdf.pages:
                        text += page.extract_text() or ""
                return text
            except Exception as e:
                print(f"Error reading PDF: {e}")
                return None
        
        # Word Document Handling
        elif attachment.filename.endswith(".docx"):
            try:
                doc = Document(file_like_object)
                text = "\n".join([para.text for para in doc.paragraphs])
                return text
            except Exception as e:
                print(f"Error reading DOCX: {e}")
                return None
        
        # PowerPoint Handling
        elif attachment.filename.endswith(".pptx"):
            try:
                ppt = Presentation(file_like_object)
                text = ""
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            except Exception as e:
                print(f"Error reading PPTX: {e}")
                return None

    except Exception as e:
        print(f"Unexpected error reading file: {e}")
        return None

# Background task that processes the question queue
async def process_queue():
    while True:
        if question_queue:
            question, interaction, is_dm = question_queue.popleft()
            gpt_response = await ask_openai(question)

            # Ensure the response does not exceed Discord's 2000 character limit
            if len(gpt_response) > 2000:
                # Split the response into chunks of 2000 characters
                for i in range(0, len(gpt_response), 2000):
                    chunk = gpt_response[i:i + 2000]
                    if is_dm:
                        await interaction.channel.send(chunk)
                    else:
                        await interaction.channel.send(f"{interaction.user.display_name}'s ChatGPT response (part):\n{chunk}")
            else:
                if is_dm:
                    await interaction.channel.send(f"Your ChatGPT response: {gpt_response}")
                else:
                    await interaction.channel.send(f"{interaction.user.display_name}'s ChatGPT response: {gpt_response}")
        
        await asyncio.sleep(1)


def check_similar_questions(new_question: str) -> str:
    threshold = 0.7

    for entry in interactions_cache:
        similarity = difflib.SequenceMatcher(None, entry['question'].lower(), new_question.lower()).ratio()
        if similarity >= threshold:
            return entry['response']
    return None

async def ask_openai(question: str):
    async with semaphore:  # Only allow a limited number of concurrent requests
        try:
            # Prepend the assistant prompt to every message
            assistant_prompt = "You are a teaching assistant in a college class designed to answer student questions. Structure your responses in a way that students can learn from these answers, like providing examples or in-depth explanations."
            full_prompt = f'{assistant_prompt}\n\nStudentQuestion: "{question}"'

            async with aiohttp.ClientSession() as session:
                async with session.post(
                    'https://api.openai.com/v1/chat/completions',
                    headers={
                        'Authorization': f'Bearer {OPENAI_API_KEY}',
                        'Content-Type': 'application/json'
                    },
                    json={
                        'model': 'gpt-4',
                        'messages': [{'role': 'user', 'content': full_prompt}]
                    }
                ) as response:
                    response_json = await response.json()
                    return response_json['choices'][0]['message']['content']
        except Exception as e:
            print(f"Error contacting OpenAI: {e}")
            return "There was an error contacting ChatGPT."

@bot.tree.command(name="new_lecture")
@commands.has_role(ADMIN_ROLE_NAME)  # Restrict to users with the "Professor" role
async def store_lecture(interaction: discord.Interaction, attachment: discord.Attachment):
    await interaction.response.send_message("Processing your lecture document...")

    # Read the file content
    file_content = await read_file_content(attachment)

    if file_content:
        filetype = attachment.filename.split('.')[-1]
        lecture_id = len(lectures_cache) + 1

        lectures_cache[lecture_id] = {
            'filename': attachment.filename,
            'filetype': filetype,
            'content': file_content
        }

        await interaction.followup.send(f"Lecture '{attachment.filename}' stored successfully. ID: {lecture_id}")
    else:
        await interaction.followup.send("Failed to process the document. Please upload a valid .pdf, .docx, or .pptx file.")


@bot.tree.command(name="summary")
async def summary(interaction: discord.Interaction, lecture_id: int):
    await interaction.response.defer()

    lecture = lectures_cache.get(lecture_id)

    if lecture:
        filename = lecture['filename']
        content = lecture['content']

        # Start the summarization process
        print("Starting text summarization...")
        summary = await summarize_text(content)

        # Truncate the summary to a maximum of 1999 characters
        if len(summary) > 1999:
            print("Summary exceeds 1999 characters, truncating...")
            summary = summary[:1999]

        print("Sending summary of lecture to Discord")
        await interaction.followup.send(f"Summary of {filename}:\n{summary}")
    else:
        await interaction.followup.send(f"No lecture found with ID: {lecture_id}")



@bot.tree.command(name="cache")
@commands.has_role(ADMIN_ROLE_NAME)  # Restrict to users with the "Professor" role
async def cache(interaction: discord.Interaction):
    await interaction.response.defer()

    lecture_info = "\n".join([f"ID: {id} - Filename: {lecture['filename']} (Type: {lecture['filetype']})"
                              for id, lecture in lectures_cache.items()]) if lectures_cache else "No lectures found."

    interaction_info = "\n".join([f"User: <@{entry['user_id']}> - Question: {entry['question']} - Response: {entry['response'][:50]}..."
                                   for entry in interactions_cache]) if interactions_cache else "No interactions found."

    response_message = "###  Lectures in Memory:\n" + lecture_info + "\n\n### Interactions in Memory:\n" + interaction_info

    await interaction.followup.send(response_message)

@bot.tree.command(name="say")
async def say(interaction: discord.Interaction, *, message: str):
    await interaction.response.send_message("Processing...", delete_after=5)
    
    # Check for similar questions in the cache
    cached_answer = check_similar_questions(message)
    if cached_answer:
        print(f"Using cached answer for question: '{message}'")
        await interaction.followup.send(f"Found a similar question in the cache: {cached_answer}")
        return
    
    # Get the response from OpenAI
    response = await ask_openai(message)
    
    # Store the new interaction in the cache
    interactions_cache.append({
        'user_id': interaction.user.id,
        'question': message,
        'response': response
    })
    
    # Send response in chunks if necessary
    if len(response) > 2000:
        for i in range(0, len(response), 2000):
            chunk = response[i:i + 2000]
            await interaction.followup.send(chunk)
    else:
        await interaction.followup.send(f"{interaction.user.display_name}'s ChatGPT response: {response[:2000]}")  # Ensure trimming for display






@bot.event
async def on_message(message: discord.Message):
    if message.author == bot.user:
        return

    if message.guild is None:
        await message.channel.send("Processing your message...", delete_after=5)
        interactions_cache.append({'user_id': str(message.author.id), 'question': message.content})
        await question_queue.put((message.content, message, True))

    await bot.process_commands(message)

@bot.tree.command(name="list")
async def list_lectures(interaction: discord.Interaction):
    await interaction.response.defer()

    if lectures_cache:
        lecture_list = "\n".join([f"ID: {id} - Filename: {lecture['filename']}" for id, lecture in lectures_cache.items()])
        await interaction.followup.send(f"### Stored Lectures:\n{lecture_list}")
    else:
        await interaction.followup.send("No lectures found in the memory.")

@bot.event
async def on_ready():
    await bot.tree.sync()
    print(f'Bot is online as {bot.user}!')

    general_channel = discord.utils.get(bot.get_all_channels(), name='general')
    if general_channel:
        await general_channel.send(
            f"Hello, class! I'm your Assistant Bot, here to help facilitate the classroom experience.\n"
            f"Here's a brief explanation of how I can assist you:\n\n"
            f"**1. `/say [message]`** - Allows you to communicate with me.\n"
            f"**2. `/list`** - Prints a list of all the documents stored in the cache.\n"
            f"**3. `/summary [ID#]`** - I will summarize the document stored in the cache based on the specified ID#.\n\n"
            f"I'm looking forward to assisting you throughout the course!"
        )

    asyncio.create_task(process_queue())

# /sayiac command that extracts text from an image (from DMs or lecture channel)
@bot.tree.command(name="sayiac")
async def sayiac(interaction: discord.Interaction, *, message: str):
    # Respond immediately with "Processing..." to end the interaction quickly then delete message to clean text channel
    await interaction.response.send_message("Processing your image and question...", delete_after=5)

    # Check if the command was sent in a DM or in a server text channel
    if interaction.guild is None:
        # The command was sent from a DM, so look for the most recent image in this DM
        asyncio.create_task(process_image_from_dm(interaction, message))
    else:
        # The command was sent from a text channel, process image from "lecture" channel
        asyncio.create_task(process_image_from_lecture(interaction, message))


# Background task to process the image in the DM
async def process_image_from_dm(interaction: discord.Interaction, user_question: str):
    # Fetch the most recent message with an image in the DM
    image_message = None
    async for msg in interaction.channel.history(limit=50):
        if msg.attachments:
            for attachment in msg.attachments:
                if attachment.filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                    image_message = msg
                    break
        if image_message:
            break

    if not image_message:
        await interaction.channel.send("No image found in your recent direct messages. Please send an image and try again.")
        return

    # Download the image and extract text using pytesseract
    image_data = await image_message.attachments[0].read()
    image = Image.open(io.BytesIO(image_data))
    extracted_text = pytesseract.image_to_string(image)

    # Construct the prompt to send to ChatGPT
    constructed_prompt = f'ImageAsText: "{extracted_text}". UserQuestion: "{user_question}".'

    # Add the constructed prompt and interaction object to the queue
    question_queue.append((constructed_prompt, interaction, True))

# Background task to process the image in the "lecture" channel
async def process_image_from_lecture(interaction: discord.Interaction, user_question: str):
    # Fetch the "lecture" channel by name
    lecture_channel = discord.utils.get(interaction.guild.text_channels, name="lecture")
    if not lecture_channel:
        await interaction.channel.send("Lecture channel not found!")
        return

    # Fetch the most recent message with an image in the "lecture" channel
    image_message = None
    async for msg in lecture_channel.history(limit=50):
        if msg.attachments:
            for attachment in msg.attachments:
                if attachment.filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                    image_message = msg
                    break
        if image_message:
            break

    if not image_message:
        await interaction.channel.send("No image found in the recent messages.")
        return

    # Download and process the image
    image_data = await image_message.attachments[0].read()
    image = Image.open(io.BytesIO(image_data))
    extracted_text = pytesseract.image_to_string(image)

    # Construct the prompt to send to ChatGPT
    constructed_prompt = f'ImageAsText: "{extracted_text}". UserQuestion: "{user_question}".'

    # Add the constructed prompt and interaction object to the queue
    question_queue.append((constructed_prompt, interaction, False))

# Error handling for app commands
@bot.tree.error
async def on_app_command_error(interaction: discord.Interaction, error):
    if isinstance(error, commands.CheckFailure):
        await interaction.response.send_message("You do not have the required role to use this command.", ephemeral=True)
    else:
        print(f"An error occurred: {error}")



# Main function to start the bot
async def main():
    try:
        await bot.start(DISCORD_TOKEN)
    except KeyboardInterrupt:
        print("Bot stopped manually")
    finally:
        await bot.close()

# Run the asynchronous main function
asyncio.run(main())